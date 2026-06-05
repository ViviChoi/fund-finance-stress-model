"""
Build the Fund Finance Stress Model slide deck.

Style: banking presentation (Helvetica, slate / white, content-dense).
Output: ../slides/Fund-Finance-Stress-Model-Deck.pptx

Re-run after changing any figure or text. Idempotent.
"""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree


# ---------------------------------------------------------------------------
# Visual system
# ---------------------------------------------------------------------------

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

INK         = RGBColor(0x1B, 0x1F, 0x29)   # body text (deep slate)
INK_SOFT    = RGBColor(0x55, 0x5B, 0x6B)   # captions / footnotes
ACCENT      = RGBColor(0xB6, 0x14, 0x2E)   # accent red — banking-style
ACCENT_SOFT = RGBColor(0xD7, 0xAE, 0xB5)
DIVIDER     = RGBColor(0xC9, 0xCD, 0xD6)
BG          = RGBColor(0xFF, 0xFF, 0xFF)
BG_PANEL    = RGBColor(0xF4, 0xF5, 0xF7)
OK          = RGBColor(0x1F, 0x7A, 0x3C)
BAD         = RGBColor(0xA8, 0x1B, 0x1B)

FONT = "Helvetica Neue"


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------

def _set_run(run, text, *, size=14, bold=False, italic=False,
             color=INK, font_name=FONT):
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def _add_textbox(slide, left, top, width, height, text,
                 *, size=14, bold=False, italic=False, color=INK,
                 align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
                 line_spacing=1.15):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    _set_run(p.add_run(), text, size=size, bold=bold, italic=italic, color=color)
    return tb


def _add_paragraph(tb, text, *, size=14, bold=False, italic=False,
                   color=INK, align=PP_ALIGN.LEFT, bullet=False,
                   space_before=4, line_spacing=1.15, indent=0):
    p = tb.text_frame.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    p.line_spacing = line_spacing
    if indent:
        p.level = indent
    if bullet:
        # Force a real bullet via XML — python-pptx textbox bullets are flaky
        pPr = p._pPr if p._pPr is not None else p._p.get_or_add_pPr()
        buChar = etree.SubElement(pPr, qn("a:buChar"))
        buChar.set("char", "•")
        buFont = etree.SubElement(pPr, qn("a:buFont"))
        buFont.set("typeface", FONT)
        prefix = "  "
        run_text = prefix + text
    else:
        run_text = text
    _set_run(p.add_run(), run_text, size=size, bold=bold, italic=italic, color=color)
    return p


def _add_rect(slide, left, top, width, height, fill_color, line_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def _add_line(slide, x1, y1, x2, y2, color=DIVIDER, weight=0.75):
    line = slide.shapes.add_connector(1, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def _slide_chrome(slide, *, page_num, total_pages, eyebrow=None, title=None):
    """Standard slide chrome — top accent bar, eyebrow + title, footer, page nbr."""
    # Top accent bar
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), ACCENT)

    # Eyebrow + title block
    if eyebrow:
        _add_textbox(
            slide, Inches(0.55), Inches(0.40), Inches(12), Inches(0.30),
            eyebrow.upper(), size=10, bold=True, color=ACCENT,
        )
    if title:
        _add_textbox(
            slide, Inches(0.55), Inches(0.70), Inches(12.2), Inches(0.75),
            title, size=24, bold=True, color=INK,
        )
        _add_line(slide, Inches(0.55), Inches(1.42),
                  Inches(12.78), Inches(1.42), color=DIVIDER, weight=0.5)

    # Footer
    _add_line(slide, Inches(0.55), Inches(7.10),
              Inches(12.78), Inches(7.10), color=DIVIDER, weight=0.5)
    _add_textbox(
        slide, Inches(0.55), Inches(7.18), Inches(10), Inches(0.30),
        "Fund Finance Stress Model  ·  Jiawen Cc  ·  Polimi EE M.Sc.  ·  "
        "June 2026",
        size=9, color=INK_SOFT,
    )
    _add_textbox(
        slide, Inches(11.5), Inches(7.18), Inches(1.3), Inches(0.30),
        f"{page_num} / {total_pages}", size=9, color=INK_SOFT,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

TOTAL_PAGES = 13


def slide_01_title(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full-bleed dark panel on left third
    _add_rect(slide, Inches(0), Inches(0), Inches(4.5), SLIDE_H, INK)
    _add_rect(slide, Inches(0), Inches(0), Inches(4.5), Inches(0.15), ACCENT)

    # Left panel content
    _add_textbox(
        slide, Inches(0.55), Inches(0.55), Inches(3.8), Inches(0.5),
        "FIRST LOOK", size=11, bold=True, color=ACCENT_SOFT,
    )
    _add_textbox(
        slide, Inches(0.55), Inches(5.6), Inches(3.8), Inches(0.4),
        "Jiawen Cc", size=14, bold=True, color=BG,
    )
    _add_textbox(
        slide, Inches(0.55), Inches(5.95), Inches(3.8), Inches(0.3),
        "EE M.Sc. candidate, Politecnico di Milano",
        size=11, color=ACCENT_SOFT,
    )
    _add_textbox(
        slide, Inches(0.55), Inches(6.25), Inches(3.8), Inches(0.3),
        "A sector-decomposed NAV facility risk model",
        size=11, italic=True, color=ACCENT_SOFT,
    )
    _add_textbox(
        slide, Inches(0.55), Inches(6.85), Inches(3.8), Inches(0.3),
        "June 2026", size=10, color=ACCENT_SOFT,
    )

    # Right side title block
    _add_textbox(
        slide, Inches(5.0), Inches(1.8), Inches(8.0), Inches(0.6),
        "A power-electronics", size=36, bold=True, color=INK,
    )
    _add_textbox(
        slide, Inches(5.0), Inches(2.5), Inches(8.0), Inches(0.6),
        "risk lens on NAV facilities", size=36, bold=True, color=INK,
    )
    _add_line(slide, Inches(5.0), Inches(3.45),
              Inches(7.0), Inches(3.45), color=ACCENT, weight=2.0)
    _add_textbox(
        slide, Inches(5.0), Inches(3.65), Inches(8.0), Inches(2.5),
        "Why \"power electronics\" is four different credits, "
        "and how sub-sector decomposition changes the covenant-breach "
        "picture for 2026 NAV / hybrid facilities.",
        size=15, italic=True, color=INK_SOFT, line_spacing=1.40,
    )


def slide_02_executive_summary(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=2, total_pages=TOTAL_PAGES,
                  eyebrow="01  ·  Executive summary",
                  title="One page, three claims")

    items = [
        ("01",
         "2026 fund finance is shifting toward NAV and hybrid facilities — "
         "and the next underwriting bottleneck is sector-specific risk "
         "modelling on concentrated portfolios."),
        ("02",
         "“Power electronics” is internally heterogeneous. EV powertrain, "
         "data-center power, renewables and industrial drives behave like "
         "four different credits — listed comparables in 2024–26 ranged "
         "from −75% to +280%."),
        ("03",
         "A small Python sensitivity model makes the gap visible: at 50% "
         "initial LTV, EV-heavy tilts breach a 65% covenant under stress; "
         "data-center heavy tilts don't. Single-bucket haircuts can't see "
         "the line."),
    ]

    top = Inches(1.85)
    for i, (num, text) in enumerate(items):
        y = top + Inches(1.55 * i)
        # number block
        _add_rect(slide, Inches(0.55), y, Inches(0.8), Inches(1.20), INK)
        _add_textbox(slide, Inches(0.55), y + Inches(0.30),
                     Inches(0.8), Inches(0.6), num,
                     size=24, bold=True, color=BG, align=PP_ALIGN.CENTER)
        # text block
        _add_textbox(slide, Inches(1.7), y + Inches(0.10),
                     Inches(11.0), Inches(1.1), text,
                     size=15, color=INK, line_spacing=1.30)


def slide_03_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=3, total_pages=TOTAL_PAGES,
                  eyebrow="02  ·  The 2026 setup",
                  title="NAV financing has moved from niche to core")

    # Left column — narrative
    tb = _add_textbox(slide, Inches(0.55), Inches(1.85),
                      Inches(7.5), Inches(0.4),
                      "What changed", size=13, bold=True, color=ACCENT)
    body = _add_textbox(slide, Inches(0.55), Inches(2.25),
                        Inches(7.5), Inches(4.5),
                        "", size=14)
    body.text_frame.text = ""
    _add_paragraph(body,
        "NAV financing is now described by the Fund Finance Association "
        "as a \"core liquidity solution across private markets\" — driven "
        "by extended exit timelines and evolving LP expectations.",
        size=13, color=INK, line_spacing=1.35, space_before=0)
    _add_paragraph(body,
        "Hybrid facilities (sub-line + NAV) are gaining ground "
        "specifically because they handle concentrated portfolios better "
        "than single-collateral structures.",
        size=13, color=INK, line_spacing=1.35, space_before=10)
    _add_paragraph(body,
        "Ropes & Gray's 2026 outlook is explicit on the growth area: "
        "\"greater sector specialisation and broader adoption across "
        "geographies and asset classes.\"",
        size=13, color=INK, line_spacing=1.35, space_before=10)

    # Right column — pullquote panel
    _add_rect(slide, Inches(8.5), Inches(1.85),
              Inches(4.3), Inches(4.5), BG_PANEL)
    _add_rect(slide, Inches(8.5), Inches(1.85),
              Inches(0.10), Inches(4.5), ACCENT)
    _add_textbox(slide, Inches(8.80), Inches(2.10),
                 Inches(4.0), Inches(0.4),
                 "THE IMPLICATION", size=10, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(8.80), Inches(2.55),
                 Inches(4.0), Inches(3.5),
                 "Lenders increasingly need underwriting capability "
                 "that can read the underlying technology cycle — "
                 "not just pull a sector haircut from a table.",
                 size=15, italic=True, color=INK, line_spacing=1.40)

    # Sources line at bottom
    _add_textbox(slide, Inches(0.55), Inches(6.65),
                 Inches(12), Inches(0.25),
                 "Sources: Macfarlanes (15th FFA Symposium); Ropes & Gray (2026 outlook); "
                 "Norton Rose Fulbright (hybrid facilities); Funds Europe (NAV financing).",
                 size=9, italic=True, color=INK_SOFT)


def slide_04_subsector_thesis(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=4, total_pages=TOTAL_PAGES,
                  eyebrow="03  ·  The sector argument",
                  title="“Power electronics” is four different credits")

    # Big claim bar
    _add_rect(slide, Inches(0.55), Inches(1.80),
              Inches(12.2), Inches(0.85), BG_PANEL)
    _add_textbox(slide, Inches(0.85), Inches(1.95),
                 Inches(11.6), Inches(0.6),
                 "A NAV facility against a fund 60% concentrated in EV "
                 "powertrain is a fundamentally different credit than "
                 "one 60% concentrated in data-center power —",
                 size=14, italic=True, color=INK, line_spacing=1.30)
    _add_textbox(slide, Inches(0.85), Inches(2.32),
                 Inches(11.6), Inches(0.3),
                 "even though both fall under the same investment-policy "
                 "category.",
                 size=14, italic=True, color=INK)

    # Table headers
    table_top = Inches(2.95)
    cols = [
        ("Sub-sector",                       Inches(0.55), Inches(3.4)),
        ("Demand driver",                    Inches(3.95), Inches(3.0)),
        ("Cyclicality",                      Inches(6.95), Inches(2.5)),
        ("Customer concentration",           Inches(9.45), Inches(3.3)),
    ]
    for label, x, w in cols:
        _add_textbox(slide, x, table_top, w, Inches(0.35),
                     label.upper(), size=10, bold=True, color=ACCENT)
    _add_line(slide, Inches(0.55), table_top + Inches(0.36),
              Inches(12.78), table_top + Inches(0.36),
              color=DIVIDER, weight=0.5)

    rows = [
        ("EV powertrain",      "Consumer auto cycle",      "High, consumer-led",
         "5–10 OEM buyers"),
        ("Data-center power",  "Hyperscaler AI capex",     "Lumpy super-cycle",
         "3–5 hyperscalers"),
        ("Renewable power conv.", "Subsidy + tariff regime", "Policy-driven",
         "Utilities + EPCs"),
        ("Industrial drives",  "Industrial CapEx / IIoT",  "Mid-cycle, defensive",
         "Fragmented"),
    ]
    row_h = Inches(0.55)
    for i, (a, b, c, d) in enumerate(rows):
        y = table_top + Inches(0.50) + row_h * i
        if i % 2 == 0:
            _add_rect(slide, Inches(0.55), y - Inches(0.05),
                      Inches(12.23), row_h, BG_PANEL)
        _add_textbox(slide, cols[0][1], y, cols[0][2], Inches(0.4),
                     a, size=12, bold=True, color=INK)
        _add_textbox(slide, cols[1][1], y, cols[1][2], Inches(0.4),
                     b, size=12, color=INK)
        _add_textbox(slide, cols[2][1], y, cols[2][2], Inches(0.4),
                     c, size=12, color=INK)
        _add_textbox(slide, cols[3][1], y, cols[3][2], Inches(0.4),
                     d, size=12, color=INK)


def slide_05_sector_matrix_figure(prs, fig_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=5, total_pages=TOTAL_PAGES,
                  eyebrow="03  ·  The sector argument (cont.)",
                  title="Sub-sectors on the axes that drive NAV facility risk")

    # Figure on the right
    slide.shapes.add_picture(str(fig_path),
                             Inches(4.2), Inches(1.75),
                             height=Inches(5.10))

    # Notes panel on the left
    _add_textbox(slide, Inches(0.55), Inches(1.85),
                 Inches(3.5), Inches(0.35),
                 "HOW TO READ", size=10, bold=True, color=ACCENT)
    tb = _add_textbox(slide, Inches(0.55), Inches(2.20),
                      Inches(3.5), Inches(4.5),
                      "", size=12)
    _add_paragraph(tb,
        "X-axis: demand cyclicality. Defensive (left) → highly "
        "cyclical (right).",
        size=11, color=INK, line_spacing=1.30, space_before=0)
    _add_paragraph(tb,
        "Y-axis: customer concentration. Fragmented (bottom) → "
        "top-3 buyers >70% of revenue (top).",
        size=11, color=INK, line_spacing=1.30, space_before=10)
    _add_paragraph(tb,
        "Bubble size: 2024–26 listed-comparable total-return "
        "dispersion (proxy for implied vol of private valuations).",
        size=11, color=INK, line_spacing=1.30, space_before=10)
    _add_paragraph(tb,
        "EV powertrain sits in the worst quadrant — cyclical AND "
        "concentrated. It is the underwriting risk most likely to "
        "be mispriced under a single-sector haircut.",
        size=11, bold=True, color=BAD, line_spacing=1.30, space_before=14)


def slide_06_empirical_reality(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=6, total_pages=TOTAL_PAGES,
                  eyebrow="04  ·  Empirical reality check",
                  title="2024–2026 listed comparables: a 350-point spread")

    rows = [
        ("Wolfspeed",          "EV powertrain (SiC)",        "≈ −70%", BAD),
        ("ON Semiconductor",   "EV powertrain",              "≈ −40%", BAD),
        ("Vertiv Holdings",    "Data-center power",          "≈ +280%", OK),
        ("Schneider Electric", "Data-center power (secure power div.)", "≈ +60%",  OK),
        ("Enphase Energy",     "Renewable power conv.",      "≈ −75%", BAD),
        ("ABB",                "Industrial drives",          "≈ +25%", OK),
    ]

    table_top = Inches(1.95)
    headers = [
        ("Listed comparable",  Inches(0.55), Inches(3.5)),
        ("Sub-sector",         Inches(4.20), Inches(5.5)),
        ("2024–26 return",     Inches(10.0), Inches(2.7)),
    ]
    for label, x, w in headers:
        _add_textbox(slide, x, table_top, w, Inches(0.35),
                     label.upper(), size=10, bold=True, color=ACCENT)
    _add_line(slide, Inches(0.55), table_top + Inches(0.36),
              Inches(12.78), table_top + Inches(0.36),
              color=DIVIDER, weight=0.5)

    row_h = Inches(0.50)
    for i, (name, sec, ret, col) in enumerate(rows):
        y = table_top + Inches(0.50) + row_h * i
        if i % 2 == 0:
            _add_rect(slide, Inches(0.55), y - Inches(0.05),
                      Inches(12.23), row_h, BG_PANEL)
        _add_textbox(slide, headers[0][1], y, headers[0][2], Inches(0.4),
                     name, size=12, bold=True, color=INK)
        _add_textbox(slide, headers[1][1], y, headers[1][2], Inches(0.4),
                     sec, size=12, color=INK)
        _add_textbox(slide, headers[2][1], y, headers[2][2], Inches(0.4),
                     ret, size=13, bold=True, color=col,
                     align=PP_ALIGN.RIGHT)

    # Pull quote
    _add_rect(slide, Inches(0.55), Inches(5.85),
              Inches(12.2), Inches(0.90), BG_PANEL)
    _add_rect(slide, Inches(0.55), Inches(5.85),
              Inches(0.10), Inches(0.90), ACCENT)
    _add_textbox(slide, Inches(0.85), Inches(5.99),
                 Inches(11.7), Inches(0.7),
                 "The dispersion is the point. Two listed names labelled "
                 "“power electronics” returned −75% and +280% in the "
                 "same window — a 350-point gap. Any underwriting model "
                 "that aggregates them is, by construction, blind.",
                 size=13, italic=True, color=INK, line_spacing=1.35)


def slide_07_model_setup(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=7, total_pages=TOTAL_PAGES,
                  eyebrow="05  ·  The model",
                  title="Sensitivity setup: same fund, two facility regimes")

    # Two side-by-side panels
    panels = [
        (Inches(0.55), "BENIGN", "Initial LTV 30%",
         ["NAV USD 500M", "Facility USD 150M", "Covenant LTV 65%"], OK),
        (Inches(7.05), "AGGRESSIVE", "Initial LTV 50%",
         ["NAV USD 500M", "Facility USD 250M", "Covenant LTV 65%"], BAD),
    ]
    panel_w = Inches(6.2)
    panel_h = Inches(2.10)
    for x, eyebrow, headline, items, col in panels:
        _add_rect(slide, x, Inches(1.85), panel_w, panel_h, BG_PANEL)
        _add_rect(slide, x, Inches(1.85), Inches(0.10), panel_h, col)
        _add_textbox(slide, x + Inches(0.30), Inches(2.00),
                     Inches(5.5), Inches(0.30),
                     eyebrow, size=10, bold=True, color=col)
        _add_textbox(slide, x + Inches(0.30), Inches(2.30),
                     Inches(5.5), Inches(0.5),
                     headline, size=18, bold=True, color=INK)
        tb = _add_textbox(slide, x + Inches(0.30), Inches(2.85),
                          Inches(5.5), Inches(1.2), "", size=12)
        for it in items:
            _add_paragraph(tb, it, size=12, bullet=True,
                           color=INK, space_before=4, line_spacing=1.30)

    # Sector stress table
    _add_textbox(slide, Inches(0.55), Inches(4.20),
                 Inches(12), Inches(0.35),
                 "SUB-SECTOR STRESS ASSUMPTIONS",
                 size=10, bold=True, color=ACCENT)
    _add_line(slide, Inches(0.55), Inches(4.56),
              Inches(12.78), Inches(4.56), color=DIVIDER, weight=0.5)

    sectors = [
        ("EV powertrain",       "−35%",
         "Calibrated to Wolfspeed / ON Semi (auto) 2024–26 drawdowns, "
         "dampened for private mark lag."),
        ("Data-center power",   "−10%",
         "Mild — secular AI-capex tailwind. Vertiv / Schneider returns positive."),
        ("Renewable power",     "−30%",
         "Policy + oversupply shock. Enphase / SolarEdge as anchors."),
        ("Industrial drives",   "−10%",
         "Mid-cycle baseline. ABB / Rockwell range stable."),
    ]
    row_h = Inches(0.42)
    for i, (name, stress, rationale) in enumerate(sectors):
        y = Inches(4.65) + row_h * i
        if i % 2 == 0:
            _add_rect(slide, Inches(0.55), y - Inches(0.04),
                      Inches(12.23), row_h, BG_PANEL)
        _add_textbox(slide, Inches(0.70), y, Inches(2.7), Inches(0.4),
                     name, size=12, bold=True, color=INK)
        _add_textbox(slide, Inches(3.40), y, Inches(1.0), Inches(0.4),
                     stress, size=13, bold=True, color=BAD,
                     align=PP_ALIGN.RIGHT)
        _add_textbox(slide, Inches(4.65), y, Inches(8.1), Inches(0.4),
                     rationale, size=11, color=INK_SOFT)


def slide_08_heatmap(prs, fig_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=8, total_pages=TOTAL_PAGES,
                  eyebrow="06  ·  Results",
                  title="Stressed LTV — same fund, two facility regimes")

    slide.shapes.add_picture(str(fig_path),
                             Inches(0.55), Inches(1.75),
                             width=Inches(12.2))

    _add_textbox(slide, Inches(0.55), Inches(6.55),
                 Inches(12), Inches(0.30),
                 "READ THE DASHED LINE", size=10, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.55), Inches(6.80),
                 Inches(12.2), Inches(0.3),
                 "Left: no tilt breaches the 65% covenant. Right: the "
                 "covenant contour cuts the feasible region — every "
                 "portfolio lower-right of the dashed line is in breach. "
                 "Same fund, same shocks; only the facility size changed.",
                 size=11, color=INK, line_spacing=1.30)


def slide_09_punchline(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=9, total_pages=TOTAL_PAGES,
                  eyebrow="07  ·  Punchline",
                  title="Reference portfolios — the dispersion in numbers")

    # Table
    headers = [
        ("Portfolio",                         Inches(0.55), Inches(4.3)),
        ("Drawdown",                          Inches(4.95), Inches(1.4)),
        ("Stressed NAV",                      Inches(6.45), Inches(1.9)),
        ("LTV @ 30% initial",                 Inches(8.45), Inches(2.1)),
        ("LTV @ 50% initial",                 Inches(10.65), Inches(2.1)),
    ]
    table_top = Inches(1.95)
    for label, x, w in headers:
        _add_textbox(slide, x, table_top, w, Inches(0.35),
                     label.upper(), size=9.5, bold=True, color=ACCENT,
                     align=PP_ALIGN.LEFT if "Portfolio" in label else PP_ALIGN.RIGHT)
    _add_line(slide, Inches(0.55), table_top + Inches(0.36),
              Inches(12.78), table_top + Inches(0.36),
              color=DIVIDER, weight=0.5)

    rows = [
        ("100% EV powertrain",       "−35.0%", "USD 325M",     "46.2%",    "76.9%   BREACH",  True),
        ("100% data-center power",   "−10.0%", "USD 450M",     "33.3%",    "55.6%",           False),
        ("Diversified (25/25/25/25)","−21.25%","USD 393.75M",  "38.1%",    "63.5%",           False),
    ]
    row_h = Inches(0.65)
    for i, (a, b, c, d, e, breach) in enumerate(rows):
        y = table_top + Inches(0.50) + row_h * i
        if i % 2 == 0:
            _add_rect(slide, Inches(0.55), y - Inches(0.05),
                      Inches(12.23), row_h, BG_PANEL)
        _add_textbox(slide, headers[0][1], y, headers[0][2], Inches(0.5),
                     a, size=13, bold=True, color=INK)
        _add_textbox(slide, headers[1][1], y, headers[1][2], Inches(0.5),
                     b, size=13, color=INK, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, headers[2][1], y, headers[2][2], Inches(0.5),
                     c, size=13, color=INK, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, headers[3][1], y, headers[3][2], Inches(0.5),
                     d, size=13, color=INK, align=PP_ALIGN.RIGHT)
        _add_textbox(slide, headers[4][1], y, headers[4][2], Inches(0.5),
                     e, size=13, bold=True,
                     color=BAD if breach else INK, align=PP_ALIGN.RIGHT)

    # Three callouts
    callouts = [
        ("01", "At 30% initial LTV, no tilt breaches. Conservative facility "
               "sizing absorbs every shock."),
        ("02", "At 50% initial LTV, EV-heavy tilts breach by ~12 points. "
               "Data-center heavy tilts stay 9 points inside."),
        ("03", "Diversified 25/25/25/25 sits 1.5pp below covenant — "
               "safe on paper, fragile in practice. A 5pp extra EV tilt "
               "breaches."),
    ]
    y0 = Inches(4.40)
    for i, (num, text) in enumerate(callouts):
        x = Inches(0.55) + Inches(4.15) * i
        _add_rect(slide, x, y0, Inches(3.95), Inches(2.20), BG_PANEL)
        _add_rect(slide, x, y0, Inches(3.95), Inches(0.05), ACCENT)
        _add_textbox(slide, x + Inches(0.20), y0 + Inches(0.20),
                     Inches(3.6), Inches(0.4),
                     num, size=22, bold=True, color=ACCENT)
        _add_textbox(slide, x + Inches(0.20), y0 + Inches(0.85),
                     Inches(3.6), Inches(1.3),
                     text, size=12, color=INK, line_spacing=1.30)


def slide_10_limits_and_next(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=10, total_pages=TOTAL_PAGES,
                  eyebrow="08  ·  Honest limits + roadmap",
                  title="What the model is not — and what comes next")

    # Two columns
    col_w = Inches(6.0)
    # Left — limits
    _add_rect(slide, Inches(0.55), Inches(1.85),
              col_w, Inches(4.9), BG_PANEL)
    _add_textbox(slide, Inches(0.80), Inches(2.00),
                 col_w, Inches(0.35),
                 "WHAT THIS MODEL IS NOT", size=10, bold=True, color=BAD)
    tb = _add_textbox(slide, Inches(0.80), Inches(2.45),
                      Inches(5.6), Inches(4.0), "", size=12)
    for item in [
        "Not a production NAV pricing model. No real advance rates, no "
        "redetermination logic, no recovery modelling.",
        "Stress numbers are illustrative anchors calibrated from listed "
        "comparables — not consensus haircuts.",
        "Cross-sub-sector correlation is implicitly zero. In practice EV "
        "and data-center are partially correlated through rates.",
        "Workout mechanics for NAV facility defaults are not modelled — "
        "this is pre-recovery LTV.",
    ]:
        _add_paragraph(tb, item, size=11.5, bullet=True,
                       color=INK, line_spacing=1.30, space_before=8)

    # Right — next steps
    _add_rect(slide, Inches(6.75), Inches(1.85),
              col_w, Inches(4.9), BG_PANEL)
    _add_textbox(slide, Inches(7.0), Inches(2.00),
                 col_w, Inches(0.35),
                 "WHAT I'D BUILD NEXT", size=10, bold=True, color=OK)
    tb = _add_textbox(slide, Inches(7.0), Inches(2.45),
                      Inches(5.6), Inches(4.0), "", size=12)
    for item in [
        "Correlation block — at minimum a two-factor decomposition: "
        "rate sensitivity and AI-capex sensitivity.",
        "Dynamic advance-rate redetermination, so a stress event "
        "triggers a borrowing-base reduction before the covenant breaches.",
        "Extend to hybrid facilities — let the undrawn subscription "
        "line act as a fallback rail when NAV LTV drifts toward covenant.",
        "Replace listed-comparable calibration with realised PE marks "
        "(your team has these; I don't).",
    ]:
        _add_paragraph(tb, item, size=11.5, bullet=True,
                       color=INK, line_spacing=1.30, space_before=8)


def slide_11_beyond(prs):
    """Beyond Power Electronics — methodology transfer (Option B framing)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=11, total_pages=TOTAL_PAGES,
                  eyebrow="09  ·  Methodology transfer",
                  title="The reflex extends beyond power electronics")

    # Lead-in panel
    _add_rect(slide, Inches(0.55), Inches(1.80),
              Inches(12.2), Inches(0.85), BG_PANEL)
    _add_textbox(slide, Inches(0.85), Inches(1.95),
                 Inches(11.6), Inches(0.30),
                 "THESIS", size=10, bold=True, color=ACCENT)
    _add_textbox(slide, Inches(0.85), Inches(2.20),
                 Inches(11.6), Inches(0.55),
                 "The decomposition reflex is not specific to power electronics. "
                 "Three sketches — sectors where the same approach applies, "
                 "deliberately without numbers to avoid claiming expertise I "
                 "don't yet have.",
                 size=12, italic=True, color=INK, line_spacing=1.35)

    # Three columns
    col_w = Inches(4.0)
    col_x_start = Inches(0.55)
    col_gap = Inches(0.15)
    col_y = Inches(2.85)
    col_h = Inches(3.30)

    columns = [
        ("Semiconductors", "AI-era semis PE fund", [
            "GPU / AI accelerators — hyperscaler concentration",
            "ASIC specialists — per-customer design cycle",
            "Advanced packaging (CoWoS / HBM) — capacity premium",
            "Lithography & deposition (ASML, Applied, Lam) — oligopoly",
            "Mature-node foundries — defensive baseline",
        ], "Customer concentration × process-node × capex-cycle"),

        ("Medical devices", "Surgical / interventional medtech PE", [
            "Surgical robotics — 7–10y replacement cycles",
            "Disposables & instruments — recurring, GMP moat",
            "Digital surgery / workflow software — SaaS-like",
            "Interventional cardiology — FDA-pipeline timing",
            "Diagnostic imaging — reimbursement-code dynamics",
        ], "FDA timing × hospital capex × reimbursement codes"),

        ("Grid modernisation", "T&D / electrification PE fund", [
            "T&D equipment — 2–4y lead-time crisis premium",
            "Grid-scale storage integrators — queue-dependent",
            "DERMS / utility software — recurring + integration",
            "EV charging infrastructure — utilisation × policy",
            "Interconnection / EPC services — queue management",
        ], "Utility capex × policy × queue length × lead-time"),
    ]

    for i, (title, subtitle, items, dispersion) in enumerate(columns):
        x = col_x_start + (col_w + col_gap) * i
        _add_rect(slide, x, col_y, col_w, col_h, BG_PANEL)
        _add_rect(slide, x, col_y, col_w, Inches(0.05), ACCENT)

        # Title + subtitle
        _add_textbox(slide, x + Inches(0.20), col_y + Inches(0.15),
                     col_w - Inches(0.30), Inches(0.35),
                     title, size=14, bold=True, color=INK)
        _add_textbox(slide, x + Inches(0.20), col_y + Inches(0.50),
                     col_w - Inches(0.30), Inches(0.25),
                     subtitle, size=9, italic=True, color=INK_SOFT)

        # Sub-sector bullets
        tb = _add_textbox(slide, x + Inches(0.20), col_y + Inches(0.85),
                          col_w - Inches(0.40), Inches(1.80),
                          "", size=10)
        for item in items:
            _add_paragraph(tb, item, size=10, bullet=True, color=INK,
                           space_before=3, line_spacing=1.25)

        # Dispersion driver line
        _add_textbox(slide, x + Inches(0.20), col_y + Inches(2.65),
                     col_w - Inches(0.30), Inches(0.20),
                     "STRESS DISPERSION FROM:", size=8,
                     bold=True, color=ACCENT)
        _add_textbox(slide, x + Inches(0.20), col_y + Inches(2.85),
                     col_w - Inches(0.30), Inches(0.40),
                     dispersion, size=9.5, italic=True, color=INK,
                     line_spacing=1.25)

    # Bottom takeaway bar
    _add_rect(slide, Inches(0.55), Inches(6.40),
              Inches(12.2), Inches(0.55), INK)
    _add_textbox(slide, Inches(0.85), Inches(6.50),
                 Inches(11.6), Inches(0.40),
                 "I can sketch the framework — calibrating any of these "
                 "to covenant-relevant numbers needs the desk's private "
                 "NAV marks. The reflex is the asset; calibration is the "
                 "conversation.",
                 size=11, italic=True, color=BG, line_spacing=1.30)


def slide_12_about(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, page_num=12, total_pages=TOTAL_PAGES,
                  eyebrow="10  ·  About me",
                  title="An engineer applying to a non-engineering seat")

    # Left — bio
    _add_textbox(slide, Inches(0.55), Inches(1.95),
                 Inches(7.0), Inches(0.4),
                 "BACKGROUND", size=10, bold=True, color=ACCENT)
    tb = _add_textbox(slide, Inches(0.55), Inches(2.35),
                      Inches(7.0), Inches(4.2), "", size=12)
    _add_paragraph(tb,
        "EE M.Sc. candidate at Politecnico di Milano, technical "
        "concentration in power electronics and control.",
        size=12.5, color=INK, line_spacing=1.35, space_before=0)
    _add_paragraph(tb,
        "Most recent project: EMI filter design for a switched-mode "
        "converter — direct hands-on exposure to the components and "
        "end-markets PE capital is currently flowing into.",
        size=12.5, color=INK, line_spacing=1.35, space_before=10)
    _add_paragraph(tb,
        "Why fund finance, why now: the analytical reflexes from EE — "
        "sensitivity, decomposition, stress testing — are the same ones "
        "NAV underwriting needs. And the sectors where NAV facilities "
        "are growing fastest in 2026 are exactly the ones where "
        "engineering literacy at the desk would differentiate the bank.",
        size=12.5, color=INK, line_spacing=1.35, space_before=10)

    # Right — what I bring / what I don't
    _add_rect(slide, Inches(8.0), Inches(1.95),
              Inches(4.8), Inches(4.55), BG_PANEL)
    _add_textbox(slide, Inches(8.25), Inches(2.10),
                 Inches(4.5), Inches(0.35),
                 "WHAT I BRING", size=10, bold=True, color=OK)
    tb1 = _add_textbox(slide, Inches(8.25), Inches(2.50),
                       Inches(4.4), Inches(1.7), "", size=11)
    for line in [
        "Engineering depth in power electronics",
        "Python + sensitivity-analysis fluency",
        "Reflex for decomposing concentrated risk",
        "Willingness to admit what I don't know",
    ]:
        _add_paragraph(tb1, line, size=11.5, bullet=True,
                       color=INK, line_spacing=1.30, space_before=4)

    _add_textbox(slide, Inches(8.25), Inches(4.50),
                 Inches(4.5), Inches(0.35),
                 "WHAT I DON'T (YET)", size=10, bold=True, color=BAD)
    tb2 = _add_textbox(slide, Inches(8.25), Inches(4.85),
                       Inches(4.4), Inches(1.6), "", size=11)
    for line in [
        "Transaction experience on NAV deals",
        "Legal / credit-document grammar",
        "Workout mechanics of facility breaches",
        "Desk-level pricing conventions",
    ]:
        _add_paragraph(tb2, line, size=11.5, bullet=True,
                       color=INK, line_spacing=1.30, space_before=4)


def slide_13_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # Full slate background
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, INK)
    _add_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.15), ACCENT)

    _add_textbox(slide, Inches(0.55), Inches(2.40),
                 Inches(12), Inches(0.5),
                 "THANK YOU", size=14, bold=True, color=ACCENT_SOFT)
    _add_textbox(slide, Inches(0.55), Inches(2.85),
                 Inches(12), Inches(0.9),
                 "Happy to walk through the model live —",
                 size=32, bold=True, color=BG)
    _add_textbox(slide, Inches(0.55), Inches(3.55),
                 Inches(12), Inches(0.9),
                 "or step back honestly if it isn't the right fit.",
                 size=32, bold=True, color=BG)

    _add_line(slide, Inches(0.55), Inches(5.40),
              Inches(2.50), Inches(5.40), color=ACCENT, weight=2.0)

    _add_textbox(slide, Inches(0.55), Inches(5.60),
                 Inches(12), Inches(0.4),
                 "Jiawen Cc  ·  Jiawen.Cc@outlook.com",
                 size=14, color=BG)
    _add_textbox(slide, Inches(0.55), Inches(6.00),
                 Inches(12), Inches(0.4),
                 "EE M.Sc. candidate  ·  Politecnico di Milano",
                 size=12, color=ACCENT_SOFT)
    _add_textbox(slide, Inches(0.55), Inches(6.40),
                 Inches(12), Inches(0.4),
                 "github.com/ViviChoi/fund-finance-stress-model",
                 size=11, italic=True, color=ACCENT_SOFT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------

def main():
    here = Path(__file__).resolve().parent
    figures = here.parent / "figures"
    slides_dir = here.parent / "slides"
    slides_dir.mkdir(exist_ok=True)
    out_path = slides_dir / "Fund-Finance-Stress-Model-Deck.pptx"

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_executive_summary(prs)
    slide_03_setup(prs)
    slide_04_subsector_thesis(prs)
    slide_05_sector_matrix_figure(prs, figures / "sector_matrix.png")
    slide_06_empirical_reality(prs)
    slide_07_model_setup(prs)
    slide_08_heatmap(prs, figures / "nav_ltv_heatmap.png")
    slide_09_punchline(prs)
    slide_10_limits_and_next(prs)
    slide_11_beyond(prs)
    slide_12_about(prs)
    slide_13_close(prs)

    prs.save(out_path)
    print(f"Deck saved to: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
